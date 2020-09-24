from pptx import Presentation
from pptx.util import Inches
from datetime import datetime, timedelta

def main():

    pres = Presentation( 'biodiversity_basic.pptx' )
    school_name = 'testing'
    
    output_textbox_numbers( 0, pres )

def output_textbox_numbers( slide_number, pres ): 
    slide = pres.slides[ slide_number ]
    
    for shape in slide.shapes: 
        if not shape.has_text_frame: 
            continue
        print( shape.text )
        
    #'{yseman}'.format( noman='noman' )

def change_slide_1( slide, school_name ):  

    for shape in slide.shapes: 
        if not shape.has_text_frame:
            continue
        print( 'before: ' + shape.text )
        
        if 'Produced for:' in shape.text:
            shape.text = 'Produced for: {}\nDate downloaded: {}'.format( school_name, datetime.now().strftime( '%d/%m/%Y' ) )
            
        else:
            print( 'no changes' )
        print( 'after: ' + shape.text )
        print() 
        
            
    for shape in slide.shapes: 
        if not shape.has_text_frame:
            continue
        print( 'remix: ' + shape.text )
    
def change_slide_11( slide ):

    slide.shapes.add_picture( 'largest_animal.jpg', Inches( 3.5 ), Inches( 2.5 ) )
    '''
    for shape in slide.shapes:
        print( '- ', end='' )
        
        if shape.has_text_frame: 
            print( shape.text )
        else:
            print( '-- {}'.format( shape.name) )
    '''
    # Get the image placeholder
    
    # Replace the image 
    
    # Save

main() 
    
    
