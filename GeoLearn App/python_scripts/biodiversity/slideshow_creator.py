import json
import os
import csv

from random import randint
from datetime import datetime, timedelta 
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image

BINOMIAL = 2
IUCN = 6
COMMON_NAME = 14

PLANT_DIET = 20
MEAT_DIET = 21
INSECT_DIET = 22
DIET = 23
MASS = 17

SINGLE_IMAGE = 320
SIX_IMAGES = 200
SIX_ONLY = 220
TOP_RIGHT_IMAGE = 115
DOBBLE_IMAGE = 120

class slideshow_creator:
    
    def __init__( self, difficulty, school_name, dir_path = os.path.dirname(os.path.realpath(__file__)) ):
        
        # Save the path as a property
        self.master_dir = dir_path + '/'
        
        print( self.master_dir )
        
        # Save the other assumed directories
        self.image_dir = dir_path + '/animal_images/' 
        
        if difficulty == 'advanced': 
            self.slideshow_path = dir_path + '/biodiversity_advanced.pptx'  
        else: # The default difficulty is basic 
            self.slideshow_path = dir_path + '/biodiversity_basic.pptx' 
        
        # Save the difficulty
        self.difficulty = difficulty
        self.school_name = school_name
        
        # Create the json template 
        self.create_template()        
        
    def create_template( self ):
         
         # Fill with the template info 
         template = ' "master_dir": "{master_fill}", \n"image_dir": "{image_fill}", \n"slideshow_path": \n"{slideshow_fill}", \n"slide_data": [ {data_fill} ]'
         
         template = template.format( master_fill = self.master_dir, image_fill = self.image_dir, slideshow_fill = self.slideshow_path, data_fill = '{}' )
         
         # save our template
         self.slideshow_json = template
         
         #return json.loads( '{' + self.slideshow_json + '}' )
    
    def create_slideshow( self ):

        if self.difficulty == 'advanced':
            self.create_advanced_slideshow_json()
            
        else: 
            self.create_basic_slideshow_json()
            self.create_slideshow_from_json()
        
    def combine_data( self, data_list ):
        
        combined_data = ""
        index = 0
        for datum in data_list:
            combined_data = combined_data + datum
            index = index + 1
            if index != len( data_list ):
                combined_data = combined_data + ', '
        return combined_data
        
    def create_animal_dictionary( self, animal_reader ):
        
        animal_dict = {}
        for animal in animal_reader:
            animal_dict[ animal[0] ] = animal
            
        return animal_dict
        
        
    def add_slide( self, slide_data, number, text=None, images=None, final=False ): 
        
        slide = ' \n"slide_number": {},'.format( str( number ) )
        
        if text != None:
            slide = slide + '\n"text_data": [ {} ]'.format( text )
            
        if images != None:
        
            if text != None: 
                slide = slide + ", " 
            slide = slide + '\n"image_data": [ {} ] '.format( images )
        
        slide_data = slide_data + self.wrap_string( slide )
        
        if not final:
            slide_data = slide_data + ', ' 
        
        return slide_data
        
        
        
    def create_text_data( self, text ):
    
        # Create a template for the text dictionary
        template = ' \n"text": ["{}"]'
        
        # Fill that template with the parameters
        text_data = template.format( text )
        
        # Wrap the string in curly braces and return it 
        return self.wrap_string( text_data )
        
    def wrap_string( self, string, comma = False ):
    
        if comma: 
            return '{' + string + '},'
        return '{' + string + '}'
        
    def create_image_data( self, image_path, x_pos, y_pos, template_size, slide_number, rotation=0.0 ): 
        
        new_path = self.alter_image_size( image_path, template_size, slide_number )
        
        left, top = self.find_image_placement( new_path, x_pos, y_pos, rotation )
        return self.create_image_data_helper( new_path, left, top, rotation )
    
    
    def create_dobble_images( self, animal_dict, image_size, slide_number ):
    
        X_POSITIONS = [ 3.3, 1.8, 4.8, 2.6, 4.1]
        Y_POSITIONS = [ 2.9, 3.6, 3.6, 5.7, 5.7 ]
        
        images = []
    
        dobble_ids = []
        
        chosen_number = 0
        chosen_index = 0
        duplicated_index = False
        
        # Choose 9 random dobble images 
        for index in range( 0, 9 ):
        
            new_id = randint( 0, 27 )
            
            while( new_id in dobble_ids ):
                new_id = randint( 0, 27 )
            
            dobble_ids.append( new_id )
            
        print( dobble_ids ) 
        
        # Choose one of the 9 to be the duplicated image
        chosen_index = randint( 0, 8 )
        print( "chosen_index: " + str( chosen_index ) )
        chosen_number = dobble_ids[ chosen_index ]
        
        # Determine if the original duplicate is on the right or left side
        if chosen_index >= 5:
            # left side
            
            duplicated_index = randint( 0, 4 )
            
        else:
            # right side
            duplicated_index = randint( 5, 9 )
            
        print( "duplicated: " + str( duplicated_index ) ) 
                      
        offset = 0
        
        # Create the set of images
        for index in range( 0, 10 ):
        
            path = ""
            x_pos = X_POSITIONS[ index % 5 ]
            y_pos = Y_POSITIONS[ index % 5 ]
            
            if index == 5:
                print( "====" )
                
            print( index - offset )
            
            
            if index >= 5:
                x_pos = x_pos + 5.8
            
            if index == duplicated_index:
                path = self.image_dir + "Dobble_" + str( chosen_number ) + '.jpg'
                offset = 1
                print( "chosen-dobble: " + str( chosen_number ) )
                
            else: 
                print( "dobble: " + str( dobble_ids[ index - offset ] ) )
                path = self.image_dir + 'Dobble_' + str( dobble_ids[ index - offset ] ) + '.jpg'
            
            images.append( self.create_image_data( path, x_pos, y_pos, image_size, slide_number, randint( -180, 180 ) ) )
           
        
        # return the combined data
        return self.combine_data( images )
    
    def prepare_image( self, image_path, x_pos, y_pos, template_size, rotation=0.0 ): 
    
        
        return self.find_image_placement( image, x_pos, y_pos, rotation )
    
    
    def create_image_data_helper( self, image_path, left, top, rotation=0.0 ): 
        
        # Create the template 
        template = '\n "image_path": "{}", "left": {}, "top": {}, "rotation": {} ' 
        
        # Fill the template
        image_data = template.format( image_path, left, top, rotation )
        
        # Wrap in curly braces and return it 
        return self.wrap_string( image_data )
        
    def alter_image_size( self, image_path, template_size, slide_number ):

        # Open the image
        image = Image.open( image_path )
        
        # Find the size of the image
        image_width, image_height = image.size

        # Take the height and width. The larger will be the one we scale with
        larger_size = float( max( image_width, image_height ) )
        
        # Find the ratio of the template size to the larger side length
        image_ratio = template_size / larger_size

        # Calculate the new size of the image based on the ratio
        image_width = round( image_width * image_ratio )
        image_height = round( image_height * image_ratio )

        # Resize the image
        resized_image = image.resize(( image_width, image_height ))

        new_path = image_path.replace( '.jpg', '_' + str( slide_number ) + '.jpg' )
        
        # Save the resized image
        resized_image.save( new_path )
        
        return new_path
        
    def find_image_placement( self, image_path, x_pos, y_pos, rotation ): 
        
        image = Image.open( image_path )
        
        pixels_per_inch = 96
        
        # Find the dimensions of the image
        image_width, image_height = image.size
        
        # Get the top-left points 
        top = (x_pos) - ( (image_width / 2) / pixels_per_inch )
        left = (y_pos)- ( (image_height / 2) / pixels_per_inch )
        
        return top, left
        
        
    def create_group_photo_images( self, animal_dict, size, slide_number, base_1, base_2=None ):
        
        if size == SIX_IMAGES:
            X_COORS = [3, 6.15, 9.3]
            Y_COORS = [3.3, 6]
        else: 
            X_COORS = [2.1, 6.3, 10.3]
            Y_COORS = [2.1, 5.9]
        
        group = []
        
        if base_2 == None:
            for index in range( 1, 7 ): 
            
                image_path = self.image_dir + base_1 + "_" + str( index ) + ".jpg"
                
                x_index = (index - 1 ) % 3
                y_index = (index - 1 ) % 2
                
                group.append( self.create_image_data( image_path, X_COORS[x_index], Y_COORS[y_index], size, slide_number ) )
        
        else:
            for index in range( 1, 4 ):
            
                image_path = self.image_dir + base_1 + "_" + str( index ) + ".jpg"
                
                #print( image_path )
                x_index = (index - 1 ) % 3
                y_index = (index - 1 ) % 2
                
                group.append( self.create_image_data( image_path, X_COORS[x_index], Y_COORS[y_index], size, slide_number ) )
            
            for index in range( 1, 4 ):
                image_path = self.image_dir + base_2 + "_" + str( index ) + ".jpg"
                x_index = (index + 2 ) % 3
                y_index = (index + 2 ) % 2
                
                group.append( self.create_image_data( image_path, X_COORS[x_index], Y_COORS[y_index], size, slide_number ) )
            

        return self.combine_data( group )
            
        
    def create_slideshow_from_json( self ):
        
        # Open the slideshow 
        pres = Presentation( self.slideshow_path )
        
        # Open the slide data dictionary and get the array 
        slides = self.slideshow_json[ 'slide_data' ]
        
        # Loop through the slides array, each being one slide's changes
        for slide_json in slides: 
        
            # Save the slide number 
            slide_number = slide_json[ 'slide_number' ] - 1
            
            # Open that slide in the presentation 
            slide = pres.slides[ slide_number ] 
            print( "slide number: " + str( slide_number + 1 ) )
            
            try: 
                text_data = slide_json[ 'text_data' ] 
                
                # Add to the slide the new next 
                for text in text_data: 
                
                    # Loop through the text shapes 
                    for shape in slide.shapes: 

                
                        # Check if the replacement is in the textbox
                        if shape.has_text_frame and '{}' in shape.text:
                        
                            for tex in shape.text_frame.paragraphs:
                                tex = tex.runs[0]
                                tex.text = tex.text.format( *text['text'] ) 
                                
                                print( tex.text )
                            print()
                        
                            break 
            except: 
                True
                
            try: 
                image_data = slide_json[ 'image_data' ]
                
                # Access the image data
                for image in image_data:
               
                        
                    # Add to the slide the new image
                    picture = slide.shapes.add_picture( image['image_path'],  Inches( image['left'] ), Inches( image['top'] ) )
                        
                    picture.rotation = image['rotation']
                        
                    print( image['image_path'] )

                    
            except:
                True   
       
        # Save it 
        pres.save( self.slideshow_path )
        
    
    def create_advanced_slideshow_json( self ):
        return True
        
    def create_basic_slideshow_json( self ):
    
        slide_data = ""
        ## Slide 1
        
        # Create the actual text that will fill the slide 
        text_1 = self.school_name
        
        text_2 = str( datetime.now().strftime( '%d/%m/%Y' ) )
        
        # Create the text dictionary with the desired dimensions 
        textdata_1 = self.create_text_data( text_1 )
        school_data = self.combine_data( [textdata_1] )
        
        textdata_2 = self.create_text_data( text_2 )
        
        slide_text = self.combine_data( [textdata_1, textdata_2] )
        
        # Add it to the total text in the slide 
        slide_data = self.add_slide( slide_data, 1, slide_text )
        
        ## Open the chosen mammals csv
        with open( self.master_dir + "/chosen_mammals_info.csv", encoding='utf8' ) as csv_file: 
        
            animal_reader = csv.reader( csv_file, delimiter=',' )
            
            animal_dict = self.create_animal_dictionary( animal_reader )
        
            ## Slide 11: Largest animal, part 1
            
            imagedata_1 = self.create_image_data( self.image_dir + 'largest_animal.jpg' , 6, 4.5, SINGLE_IMAGE, 11 )
        
            slide_text = school_data
            slide_images = self.combine_data( [imagedata_1] )
        
            slide_data = self.add_slide( slide_data, 11, slide_text, slide_images )
            
            ## Slide 12: Largest animal, part 2
            
            slide_text = self.create_basic_animal_text( animal_dict, 'largest_animal' )
            
            image = self.create_image_data( self.image_dir + 'largest_animal.jpg', 11, 1.25, TOP_RIGHT_IMAGE, 12 )
            
            
            slide_images = self.combine_data( [image] )
            
            slide_data = self.add_slide( slide_data, 12, slide_text, slide_images )

            ## Slide 13: Second Largest animal, part 1
            slide_number = 13
            
            # Text
            slide_text = school_data
            
            # Image
            image = self.create_image_data( self.image_dir + 'second_largest_animal.jpg', 6, 4.5, SINGLE_IMAGE, slide_number )
            
            slide_images = self.combine_data( [image] )
            
            slide_data = self.add_slide( slide_data, slide_number, slide_text, slide_images )

            ## Slide 14: Second Largest animal, part 2
            slide_number = 14
            
            # Text
            slide_text = self.create_basic_animal_text( animal_dict, 'second_largest_animal' )
            
            ## Image
            image = self.create_image_data( self.image_dir + 'second_largest_animal.jpg', 11, 1.25, TOP_RIGHT_IMAGE, slide_number )
            slide_images = self.combine_data( [image] )
            
            slide_data = self.add_slide( slide_data, 14, slide_text, slide_images )
            
            ### Slide 15: Largest Predator, part 1
            slide_number = 15
            
            ## Image
            image = self.create_image_data( self.image_dir + 'largest_predator.jpg', 6, 4.5, SINGLE_IMAGE, slide_number )
            slide_images = self.combine_data( [image] )
            
            slide_data = self.add_slide( slide_data, slide_number, images=slide_images )
            
            ### Slide 16: Largest Predator, part 2
            slide_number = 16
            
            ## Text
            
            slide_text = self.create_basic_animal_text( animal_dict, 'largest_predator' )
            
            ## Image
            image = self.create_image_data( self.image_dir + 'largest_predator.jpg', 11, 1.25, TOP_RIGHT_IMAGE, slide_number )
            slide_images = self.combine_data( [image] )
            
            slide_data = self.add_slide( slide_data, 16, slide_text, slide_images )
            
            ### Slide 17: Second Largest Predator, part 1
            slide_number = 17
            
            ## Image
            image = self.create_image_data( self.image_dir + 'second_largest_predator.jpg', 6, 4.5, SINGLE_IMAGE, slide_number )
            slide_images = self.combine_data( [image] )
            
            slide_data = self.add_slide( slide_data, slide_number, images=slide_images )
            
            ### Slide 18: Second Largest Predator, part 2
            slide_number = 18
            
            ## Text
            slide_text = self.create_basic_animal_text( animal_dict, 'second_largest_predator' )
            
            ## Image
            image = self.create_image_data( self.image_dir + 'second_largest_predator.jpg', 11, 1.25, TOP_RIGHT_IMAGE, slide_number )
            slide_images = self.combine_data( [image] )
            
            slide_data = self.add_slide( slide_data, 18, slide_text, slide_images )
            
            ### Slide 19: 6 Other large animals, part 1
            slide_number = 19
            
            ## Text            
            slide_text = school_data
            
            ## Images
            slide_images = self.create_group_photo_images( animal_dict, SIX_IMAGES, slide_number, 'large_herbivores', 'large_predators' )
            
            slide_data = self.add_slide( slide_data, 19, slide_text, slide_images )
            
            ### Slide 20: 6 Other large animals, part 2
            slide_number = 20
            
            ## Text
            
            slide_text = self.create_group_animal_names( animal_dict, 'large_herbivores', 'large_predators' )
            
            ## Images
            slide_images = self.create_group_photo_images( animal_dict, SIX_ONLY, slide_number, 'large_herbivores', 'large_predators' )
            
            slide_data = self.add_slide( slide_data, 20, slide_text, slide_images )
            
            ### Slide 21: largest past animal, part 1
            slide_number = 21
            
            ## Text
            slide_text = school_data
            
            ## Image
            image = self.create_image_data( self.image_dir + 'largest_past_animal.jpg', 6, 4.5, SINGLE_IMAGE, slide_number )
            slide_images = self.combine_data( [image] )
            
            slide_data = self.add_slide( slide_data, slide_number, slide_text, slide_images )
            
            # Largest past animal
            slide_data = self.add_slide( slide_data, 21, slide_text )
            
            ### Slide 22: largest past animal, part 2
            slide_number = 22
            
            ## text
            slide_text = self.create_basic_animal_text( animal_dict, 'largest_past_animal', historic=True )
            
            ## image
            image = self.create_image_data( self.image_dir + 'largest_past_animal.jpg', 11, 1.25, TOP_RIGHT_IMAGE, slide_number )
            slide_images = self.combine_data( [image] )
            
            # largest past animal 
            
            slide_data = self.add_slide( slide_data, 22, slide_text, slide_images )
            
            ### slide 23: largest past predator, part 1
            slide_number = 23
            
            slide_text = school_data
            
            ## image
            image = self.create_image_data( self.image_dir + 'largest_historic_predator.jpg', 6, 4.5, SINGLE_IMAGE, slide_number )
            slide_images = self.combine_data( [image] )
            # largest past predator
            
            slide_data = self.add_slide( slide_data, slide_number, slide_text, slide_images )
            
            ### slide 24: largest past predator, part 2
            slide_number = 24
            
            ## text
            slide_text = self.create_basic_animal_text( animal_dict, 'largest_historic_predator', historic=True )
            
            ## image
            image = self.create_image_data( self.image_dir + 'largest_historic_predator.jpg', 11, 1.25, TOP_RIGHT_IMAGE, slide_number )
            slide_images = self.combine_data( [image] )
            
            slide_data = self.add_slide( slide_data, slide_number, slide_text, slide_images )
            
            ### slide 25: 6 past animals, part 1
            slide_number = 25
            
            ## text
            
            # school name
            slide_text = school_data
            
            ## image
            slide_images = self.create_group_photo_images( animal_dict, SIX_IMAGES, slide_number, 'historic' )
            
            slide_data = self.add_slide( slide_data, 25, slide_text, slide_images )
            
            ### slide 26: 6 past animals, part 2
            slide_number = 26
            
            ## text
            
            slide_text = self.create_group_animal_names( animal_dict, 'historic' )
            
            ## images
            slide_images = self.create_group_photo_images( animal_dict, SIX_ONLY, slide_number, 'historic' )
            
            # 6 past animal images
            slide_data = self.add_slide( slide_data, 26, slide_text, slide_images )
            
            ## Slide 32
            slide_number = 32
            
            slide_images = self.create_dobble_images( animal_dict, DOBBLE_IMAGE, slide_number )
            
            slide_data = self.add_slide( slide_data, slide_number, images=slide_images )
            
            ## Slide 33
            slide_number = 33
            slide_images = self.create_dobble_images( animal_dict, DOBBLE_IMAGE, slide_number )
            
            slide_data = self.add_slide( slide_data, slide_number, images=slide_images )
            
            ## Slide 34
            slide_number = 34
            slide_images = self.create_dobble_images( animal_dict, DOBBLE_IMAGE, slide_number )
            
            slide_data = self.add_slide( slide_data, slide_number, images=slide_images )
            
            ## Slide 35
            slide_number = 35
            slide_images = self.create_dobble_images( animal_dict, DOBBLE_IMAGE, slide_number )
            
            slide_data = self.add_slide( slide_data, slide_number, images=slide_images )
            
            ## Slide 36
            slide_number = 36
            
            slide_images = self.create_dobble_images( animal_dict, DOBBLE_IMAGE, slide_number )
            
            slide_data = self.add_slide( slide_data, slide_number, images=slide_images, final=True )
            
            
            
            
        
        # Save the presentation 
        slideshow_string = '{' + self.slideshow_json.format( slide_data ) + '}'  
        
        print()
        index = 1
        for line in slideshow_string.split( '\n' ): 
            print( str( index ) + ": " + line )
            index = index + 1
        #print( slideshow_string )
        print()
        self.slideshow_json = json.loads( slideshow_string )
        
    
        
        
    def create_basic_animal_text( self, animal_dict, animal_title, historic=False ):
        
        name = self.create_text_data( animal_dict[animal_title][COMMON_NAME] )
        mass = self.create_text_data( str( float( animal_dict[animal_title][MASS])/1000) )
        plant_diet = self.create_text_data( animal_dict[animal_title][PLANT_DIET] )
        meat_diet = self.create_text_data( animal_dict[animal_title][MEAT_DIET] )
        insect_diet = self.create_text_data( animal_dict[animal_title][INSECT_DIET] )
        
        if not historic: 
            iucn = self.create_text_data( animal_dict[animal_title][IUCN] )
            return self.combine_data([ name, mass, plant_diet, iucn, meat_diet, insect_diet ])
           
        
        return self.combine_data([ name, mass, plant_diet, meat_diet, insect_diet ])
       
            
    def create_group_animal_names( self, animal_dict, base_1, base_2=None ): 
    
        slide_text = []
        
        if base_2 == None: 
            for index in range( 1, 7 ):
        
                animal_title = base_1 + "_" + str( index )
                textdata = self.create_text_data(animal_dict[animal_title][BINOMIAL])     
                slide_text.append( textdata )
                
        else: 
            for index in range( 1, 4 ): 
            
                animal_title = base_1 + "_" + str( index )
                textdata = self.create_text_data( animal_dict[ animal_title ][COMMON_NAME] )
                slide_text.append( textdata )
                
            for index in range( 1, 4 ): 
                animal_title = base_2 + '_' + str( index ) 
                textdata = self.create_text_data( animal_dict[ animal_title ][COMMON_NAME] )
                slide_text.append( textdata ) 
            
        return self.combine_data( slide_text )
    
    
        
        
        
        
        
        
        
